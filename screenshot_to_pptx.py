import os
import sys
import time
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

def capture_slides_to_images():
    print("Initializing Playwright...")
    os.makedirs('images_pptx', exist_ok=True)
    
    with sync_playwright() as p:
        print("Launching headless Chromium browser...")
        browser = p.chromium.launch(headless=True)
        # Create context at 1920x1080 (standard 16:9 HD resolution)
        context = browser.new_context(viewport={"width": 1920, "height": 1080})
        page = context.new_page()
        
        html_path = os.path.abspath("index.html")
        url = f"file://{html_path}"
        print(f"Loading local HTML file: {url}")
        
        page.goto(url)
        page.wait_for_selector(".slide")
        
        # Inject custom styles for clean screenshots:
        # 1. Disable all transitions & animations so things are immediately at their final states
        # 2. Hide slide counter and progress bar for perfect standalone slides
        # 3. For Slide 10 (the video slide), make sure the video wrapper transitions are instant too
        style_content = """
            *, *::before, *::after {
                animation-duration: 0s !important;
                animation-delay: 0s !important;
                transition-duration: 0s !important;
                transition-delay: 0s !important;
            }
            .slide-counter, .progress-container {
                display: none !important;
            }
        """
        page.evaluate(f"""
            const style = document.createElement('style');
            style.textContent = `{style_content}`;
            document.head.appendChild(style);
        """)
        
        total_slides = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Found {total_slides} slides in index.html")
        
        video_box = None
        
        for i in range(total_slides):
            print(f"Processing Slide {i+1} of {total_slides}...")
            # Make only slide i active
            page.evaluate(f"""
                const slides = document.querySelectorAll('.slide');
                slides.forEach((slide, idx) => {{
                    slide.classList.toggle('active', idx === {i});
                }});
            """)
            
            # Wait for layout/rendering to settle
            time.sleep(0.4)
            
            # Special handling for slide 10 (index 9) which contains the video
            if i == 9:
                print("Slide 10 detected (Video Demo). Seeking video to 1.5 seconds for poster frame...")
                has_video = page.evaluate("""() => {
                    const video = document.getElementById('demoVideo');
                    if (video) {
                        video.muted = true;
                        video.currentTime = 1.5; // Seek to 1.5 seconds when looker agent output is displayed
                        return true;
                    }
                    return false;
                }""")
                if has_video:
                    time.sleep(1.2) # Wait for seek to complete and render
                    video_box = page.locator('#demoVideoWrapper').bounding_box()
                    print(f"Calculated video box in pixels: {video_box}")
            
            screenshot_path = f"images_pptx/slide_{i}.png"
            page.screenshot(path=screenshot_path)
            print(f"Captured screenshot: {screenshot_path}")
            
        browser.close()
        print("Playwright screenshot capture finished successfully.")
        return video_box, total_slides

def set_video_autoplay(slide):
    """Edits the slide XML to make any embedded media play automatically on slide entrance."""
    conds = slide._element.xpath('.//p:timing//p:cond')
    changed = False
    for cond in conds:
        if cond.get('delay') == 'indefinite':
            print("XML Hack: Setting video start condition delay to '0' for autoplay")
            cond.set('delay', '0')
            changed = True
    return changed

def compile_pptx_presentation(video_box, total_slides):
    print("\nInitializing python-pptx compiler...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen width
    prs.slide_height = Inches(7.5)    # 16:9 widescreen height
    
    # We will use blank slide layouts
    blank_layout = prs.slide_layouts[6]
    
    for i in range(total_slides):
        print(f"Assembling Slide {i+1} into PPTX...")
        slide = prs.slides.add_slide(blank_layout)
        
        screenshot_path = f"images_pptx/slide_{i}.png"
        if not os.path.exists(screenshot_path):
            print(f"Error: Screenshot missing for slide {i}, skipping.")
            continue
            
        # Add the 16:9 screenshot as full-bleed background (fits perfectly, no stretching!)
        slide.shapes.add_picture(
            screenshot_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        # If this is Slide 10 (index 9, video slide), insert the native playable video overlay
        if i == 9 and video_box is not None:
            print("Embedding native video on Slide 10...")
            video_path = os.path.join('media', 'looker-agent-cropped.mp4')
            if os.path.exists(video_path):
                # 1. Calculate matching position in Inches from Playwright pixel coordinates
                # Playwright operates at 1920x1080 viewport. PowerPoint is 13.333" x 7.5"
                px_to_inch_x = 13.333 / 1920
                px_to_inch_y = 7.5 / 1080
                
                left_inch = video_box['x'] * px_to_inch_x
                top_inch = video_box['y'] * px_to_inch_y
                width_inch = video_box['width'] * px_to_inch_x
                height_inch = video_box['height'] * px_to_inch_y
                
                print(f"Video physical alignment in PPTX: Left={left_inch:.3f}\", Top={top_inch:.3f}\", Width={width_inch:.3f}\", Height={height_inch:.3f}\"")
                
                # 2. Extract a perfect poster frame from the slide screenshot corresponding to the video box
                try:
                    print("Generating high-fidelity poster frame from cropped screenshot...")
                    img = Image.open(screenshot_path)
                    crop_area = (
                        int(video_box['x']),
                        int(video_box['y']),
                        int(video_box['x'] + video_box['width']),
                        int(video_box['y'] + video_box['height'])
                    )
                    cropped_img = img.crop(crop_area)
                    poster_path = "images_pptx/video_poster.png"
                    cropped_img.save(poster_path)
                    print(f"Poster frame saved successfully: {poster_path}")
                except Exception as e:
                    print(f"Warning: Failed to crop poster frame: {e}. Falling back to none.")
                    poster_path = None
                
                # 3. Add movie shape overlay with poster frame
                slide.shapes.add_movie(
                    video_path,
                    left=Inches(left_inch),
                    top=Inches(top_inch),
                    width=Inches(width_inch),
                    height=Inches(height_inch),
                    poster_frame_image=poster_path,
                    mime_type='video/mp4'
                )
                print("Video successfully embedded and overlaid over background screenshot.")
                
                # 4. Modify the underlying slide XML timing tree to autoplay the video
                set_video_autoplay(slide)
            else:
                print(f"Warning: Video file '{video_path}' not found. Skipping embedding.")
                
    output_filename = "Escape_Velocity_Presentation.pptx"
    prs.save(output_filename)
    print(f"\nSuccess! High-fidelity presentation compiled successfully to '{output_filename}' ({os.path.getsize(output_filename)} bytes)")

if __name__ == "__main__":
    try:
        video_box, total_slides = capture_slides_to_images()
        compile_pptx_presentation(video_box, total_slides)
    except Exception as e:
        print(f"Execution failed: {e}", file=sys.stderr)
        sys.exit(1)
